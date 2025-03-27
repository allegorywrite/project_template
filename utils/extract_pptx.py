#!/usr/bin/env python3
"""
PPTXファイルからテキストを抽出するスクリプト
"""

import sys
import os
from pathlib import Path
from pptx import Presentation

def extract_text_from_pptx(pptx_path, output_dir=None):
    """
    PPTXファイルからテキストを抽出する
    
    Args:
        pptx_path (str or Path): PPTXファイルのパス
        output_dir (str or Path, optional): 出力ディレクトリ
    """
    pptx_path = Path(pptx_path)
    
    if not pptx_path.exists():
        print(f"エラー: ファイルが見つかりません: {pptx_path}")
        return
    
    if output_dir is None:
        output_dir = pptx_path.parent
    else:
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
    
    # 出力ファイル名
    markdown_file = output_dir / f"{pptx_path.stem}_text.md"
    
    try:
        # PPTXファイルを開く
        prs = Presentation(pptx_path)
        
        # テキストを抽出
        full_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = [f"## Slide {i+1}"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            full_text.append("\n\n".join(slide_text))
        
        # 結果を保存
        print(f"結果を保存しています: {markdown_file}")
        with open(markdown_file, "w", encoding="utf-8") as f:
            f.write("\n\n".join(full_text))
        print(f"結果を保存しました: {markdown_file}")
        
        return "\n\n".join(full_text)
        
    except Exception as e:
        print(f"エラー: {e}")
        import traceback
        traceback.print_exc()
        return None

def main():
    if len(sys.argv) < 2:
        print("使用方法: python extract_pptx.py <pptx_path> [output_dir]")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    extract_text_from_pptx(pptx_path, output_dir)

if __name__ == "__main__":
    main()
